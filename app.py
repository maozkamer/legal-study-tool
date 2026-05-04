import os
import io
import uuid
import json
import logging
from datetime import date, datetime
from flask import Flask, request, render_template, jsonify, send_file, send_from_directory
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import PyPDF2
from pptx import Presentation
from docx import Document
from docx.oxml import OxmlElement
import anthropic

logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
log = logging.getLogger(__name__)

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50 MB

limiter = Limiter(app=app, key_func=get_remote_address, default_limits=["100 per hour", "20 per minute"])

client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

CLAUDE_MODEL = "claude-opus-4-5"

# In-memory text cache: token → extracted text (up to 50 entries)
_TEXT_CACHE: dict[str, str] = {}
_MAX_CACHE  = 50


@app.errorhandler(413)
def too_large(e):
    return jsonify({"error": "הקובץ גדול מדי — מקסימום 50MB"}), 413


# ─────────────────────────────────────────────────────────────
#  Text extraction
# ─────────────────────────────────────────────────────────────

def extract_pdf(data: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(data))
    return "\n".join(p.extract_text() or "" for p in reader.pages)


def extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(r.text for r in para.runs).strip()
                    if line:
                        lines.append(line)
    return "\n".join(lines)


def extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def extract_txt(data: bytes) -> str:
    for enc in ("utf-8", "windows-1255", "iso-8859-8", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


# ─────────────────────────────────────────────────────────────
#  Document type detection
# ─────────────────────────────────────────────────────────────

def detect_type(text: str) -> str:
    sample = text[:3000]
    verdict_hits     = sum(1 for w in ["בית משפט", "נגד", "השופט", "הנאשם", "המשיב", "פסק דין"] if w in sample)
    legislation_hits = sum(1 for w in ["חוק", "סעיף", "תקנות", "כנסת", "ספר החוקים"] if w in sample)
    if verdict_hits >= 2:
        return "verdict"
    if legislation_hits >= 2:
        return "legislation"
    return "study"


# ─────────────────────────────────────────────────────────────
#  Prompts
# ─────────────────────────────────────────────────────────────

PROMPTS = {
    "verdict": """\
אתה עוזר לימודים משפטי. סכם את פסק הדין הבא בעברית בפורמט הזה בדיוק:

📋 **פרטי התיק:** [בית משפט, שנה, שמות הצדדים]
⚖️ **השאלה המשפטית:** [מה נדון]
📖 **עובדות המקרה:** [עובדות עיקריות]
💬 **טענות הצדדים:** [תובע מול נתבע]
🔨 **פסיקה:** [מה פסק בית המשפט ומדוע]
📌 **העיקרון המשפטי לבחינה:** [ההלכה החשובה]

המסמך:
""",
    "legislation": """\
אתה עוזר לימודים משפטי. סכם את החקיקה הבאה בעברית בפורמט הזה בדיוק:

📋 **שם החוק ושנה:** [שם מלא + שנת חקיקה]
🎯 **מטרת החוק:** [מה החוק בא להסדיר]
📖 **סעיפים מרכזיים:** [הסעיפים החשובים]
⚠️ **חריגים חשובים:** [חריגים, הגנות, תנאים]
📌 **מה חשוב לדעת לבחינה:** [נקודות לשינון]

המסמך:
""",
    "study": """\
אתה עוזר לימודים משפטי. סכם את חומר הלימוד הבא בעברית בפורמט הזה בדיוק:

🎯 **נושא המסמך:** [נושא ראשי]
📖 **רעיונות מרכזיים:**
1. [רעיון ראשון]
2. [רעיון שני]
3. [המשך לפי הצורך]
💡 **מושגים חשובים:** [הגדרות ומושגי מפתח]
❓ **שאלות אפשריות לבחינה:**
- [שאלה 1]
- [שאלה 2]
- [שאלה 3]
📌 **סיכום:** [תמצית בשורה אחת]

המסמך:
""",
}

TYPE_LABELS = {
    "verdict":     ("⚖️", "פסק דין"),
    "legislation": ("📜", "חקיקה"),
    "study":       ("📚", "חומר לימוד"),
}


# ─────────────────────────────────────────────────────────────
#  Claude helper
# ─────────────────────────────────────────────────────────────

def _claude(prompt: str, max_tokens: int = 2048) -> str:
    response = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=max_tokens,
        messages=[{"role": "user", "content": prompt}],
    )
    return response.content[0].text


# ─────────────────────────────────────────────────────────────
#  Long-transcript chunking helpers
# ─────────────────────────────────────────────────────────────

CHUNK_SIZE    = 40_000   # characters per chunk
MAX_CHUNKS    = 4        # max 4 chunks (~3 hours of recording)
_CHUNK_OVERLAP = 500     # character overlap between consecutive chunks

# Shared case-law detection instructions injected into every full-summary prompt
_CASE_LAW_INSTRUCTIONS = (
    "זיהוי פסקי דין:\n"
    '- חפש כל אזכור של שם תיק בפורמט: X נגד Y, X v Y, ע"א XXXX, בג"ץ XXXX, רע"א XXXX\n'
    "- גם אם המרצה אומר רק שם משפחה כמו 'פרשת כהן' או 'עניין לוי' — זה פסק דין\n"
    "- חלץ: שם התיק, העיקרון המשפטי שהמרצה הסביר, הרלוונטיות לנושא השיעור\n"
    "- אם לא הוזכרו פסקי דין — case_law=[]\n\n"
)

# Full JSON schema string reused in both direct and merge prompts
_JSON_SCHEMA = (
    "{\n"
    '  "subject": "שם הנושא המשפטי המדויק (דיני עונשין / דיני חוזים / וכו\')",\n'
    '  "table_of_contents": ["נושא 1", "נושא 2"],\n'
    '  "sections": [\n'
    '    {"level": 1, "heading": "כותרת ראשית", "content": "תוכן מלא ומקיף עם דוגמאות, ניתוחים, הרחבות, סיכום ביניים ושאלות עצמיות"},\n'
    '    {"level": 2, "heading": "תת-נושא", "content": "..."}\n'
    "  ],\n"
    '  "concepts": [{"term": "מושג", "definition": "הגדרה מדויקת ומקצועית", "example": "דוגמה ספציפית"}],\n'
    '  "comparison_tables": [\n'
    '    {"title": "השוואה בין X ל-Y", "columns": ["קריטריון", "X", "Y"],\n'
    '     "rows": [["מאפיין 1", "ערך X", "ערך Y"], ["מאפיין 2", "ערך X", "ערך Y"]]}\n'
    "  ],\n"
    '  "case_law": [\n'
    '    {"name": "שם פסק הדין", "facts": "עובדות המקרה",\n'
    '     "arguments": "טענות הצדדים", "ruling": "הכרעת בית המשפט",\n'
    '     "principle": "ההלכה / העיקרון", "relevance": "רלוונטיות לנושא"}\n'
    "  ],\n"
    '  "statutes": [\n'
    '    {"law": "שם החוק", "section": "סעיף X",\n'
    '     "text": "לשון הסעיף המלאה", "analysis": "ניתוח רכיבי הסעיף (יסוד עובדתי, יסוד נפשי, נסיבות, תוצאה)"}\n'
    "  ],\n"
    '  "scales": [\n'
    '    {"title": "סולם X (מהחמור לקל)",\n'
    '     "levels": [{"name": "רמה 1", "description": "תיאור"}, {"name": "רמה 2", "description": "תיאור"}]}\n'
    "  ],\n"
    '  "important_moments": ["רגעים חשובים שהמרצה הדגיש — מסימוני ⭐"],\n'
    '  "exam_warnings": ["⚠️ נושא שייבחן", "❗ נקודת בלבול נפוצה"],\n'
    '  "self_check_questions": ["🤔 שאלה לחזרה 1", "🤔 שאלה 2"],\n'
    '  "related_topics": "קישורים לנושאים אחרים בקורס",\n'
    '  "instructor_emphasis": ["נקודות שהמרצה הדגיש"],\n'
    '  "key_points": ["נקודה 1", "נקודה 2", "נקודה 3", "נקודה 4", "נקודה 5"],\n'
    '  "transcript_unclear_zones": ["אזורים בתמלול שלא היו ברורים — בדוק בהקלטה"]\n'
    "}"
)

# Academic-level instructions for the comprehensive lecture summary prompt
_ACADEMIC_INSTRUCTIONS = (
    "=== עקרונות הסיכום (חובה לעקוב אחר כולם) ===\n\n"
    "1. שפה: עברית בלבד. מונחים מקצועיים באנגלית מותרים (Mens Rea, corpus delicti וכו'). "
    "כתיבה זורמת, דידקטית, בסגנון מרצה שמדבר אל הסטודנט — לא ניסוח רובוטי.\n\n"
    "2. אורך: ככל שצריך. שיעור של שעה → 6–12 עמודי סיכום מלא. "
    "אל תוותר על דוגמאות, ניתוחים, או הסברים בגלל אורך.\n\n"
    "3. מבנה: level 1 = נושא ראשי, level 2 = תת-נושא. הצג היררכיה ברורה.\n\n"
    "4. סגנון דיאלוגי: \"השאלה שלנו היא...\", \"נשאל...\", \"בואו נבחן...\", "
    "\"מה זה אומר בפועל?\". שאלות מנחות שמובילות את הקורא דרך הלוגיקה.\n\n"
    "5. דוגמאות: כל מושג תיאורטי ← דוגמה ספציפית וממחישה. "
    "סמן: \"לדוג':\", \"מקרה:\", \"דוגמה:\". "
    "אם המרצה נתן דוגמה — השתמש בה. אם לא — הבא דוגמה משפטית מוכרת מהידע שלך.\n\n"
    "6. סעיפי חוק: צטט לשון מלאה (מהידע שלך אם יודע) → שדה text. "
    "נתח לרכיבים (יסוד עובדתי, יסוד נפשי, נסיבות, תוצאה) → שדה analysis. "
    "סמן: \"📜 לשון הסעיף:\".\n\n"
    "7. פסקי דין — ניתוח אקדמי מלא: שם, עובדות, טענות, הכרעה, הלכה, רלוונטיות. "
    "אם המרצה הזכיר רק שם — הרחב מהידע שלך אם מכיר. סמן: \"⚖️\".\n\n"
    "8. מינוח: תרגם לטרמינולוגיה מקצועית גם אם המרצה דיבר בשפה פשוטה "
    "(רשלנות, אחריות קפידה, קוגנטי, קש\"ס עובדתי, השתק שיפוטי וכו').\n\n"
    "9. הרחבות Claude: \"💭 הערת הרחבה:\" — ציין בבירור מה Claude מוסיף. "
    "הקורא חייב לדעת מה אמר המרצה ומה הוסיף Claude.\n\n"
    "10. הצלבות: \"(ראה גם: [נושא])\" כשרלוונטי לחומר אחר בקורס.\n\n"
    "11. סולמות: מושגים מדורגים (מחשבה פלילית < רשלנות < אחריות קפידה) → "
    "שדה scales עם levels מהחמור לקל.\n\n"
    "12. השוואות: מושגים שמתבלבלים ביניהם → comparison_tables עם columns ו-rows.\n\n"
    "13. אזהרות: ⚠️ ו-❗ → exam_warnings. "
    "נושאים שהמרצה הדגיש, סימוני ⭐ בתמלול, \"שימו לב\", \"זה ייבחן\".\n\n"
    "14. סיכום ביניים: \"📌 סיכום ביניים:\" + 3–4 נקודות — בסוף כל section level 1, בתוך content.\n\n"
    "15. שאלות עצמיות: 2–3 שאלות פתוחות בלי תשובות → self_check_questions.\n\n"
    "16. סמנים: 📜 לשון החוק | ⚖️ פסיקה | 🎓 הסבר המרצה | 💭 הרחבת Claude.\n\n"
    "17. אי-בהירות: כשהתמלול לא ברור (רעש/חיתוך) → transcript_unclear_zones. "
    "אל תנחש ואל תמציא.\n\n"
    "18. תוכן עניינים: רשימה ממוספרת של כל הנושאים הראשיים → table_of_contents.\n\n"
    "=== זיהוי פסקי דין וחוקים ===\n\n"
    + _CASE_LAW_INSTRUCTIONS
    + "זיהוי חוקים: כל \"סעיף X לחוק Y\" — חלץ, צטט לשון מלאה (אם יודע), נתח.\n"
    "אם החוק לא מוכר — ציין רק מה שהמרצה אמר עליו.\n"
    "אם אין פסקי דין — case_law=[]. אם אין חוקים — statutes=[].\n"
    "אם אין סולמות — scales=[]. אם אין השוואות — comparison_tables=[].\n"
    "אם אין אזורים לא ברורים — transcript_unclear_zones=[].\n\n"
    "החזר JSON בלבד — ללא markdown, ללא טקסט לפני או אחרי הסוגריים:\n"
)


def _full_lecture_prompt(lesson_name: str, today: str, now_time: str, duration: str, transcript_section: str) -> str:
    return (
        "אתה עוזר לימודים משפטי ברמה אקדמית גבוהה. "
        "אתה מסכם שיעורי משפטים בסגנון של מרצה אקדמי מצטיין — מעמיק, דידקטי, מקצועי. "
        "הסיכום משמש את הסטודנט כתחליף מלא לספר/לשיעור, ולכן חייב להיות מקיף ואיכותי.\n\n"
        f'פרטי השיעור:\nשם: "{lesson_name}"\nתאריך: {today}, שעה: {now_time}, משך: {duration}\n\n'
        "=== עקרונות הסיכום (חובה לעקוב אחר כולם) ===\n\n"
        "1. עברית בלבד: הסיכום כולו בעברית, ללא מילים באנגלית פרט למונחים מקצועיים מוכרים בלבד "
        "(Mens Rea, corpus delicti). כתיבה זורמת בסגנון של מרצה טוב — לא ניסוח רובוטי.\n\n"
        "2. אורך: ככל שצריך. אל תקצר. שיעור של שעה יכול להפיק 6-12 עמודי סיכום. "
        "אל תוותר על דוגמאות, ניתוחים, או הסברים בגלל אורך.\n\n"
        "3. מבנה היררכי: זהה את הנושא הראשי ופרק לעקרונות / תתי-נושאים. "
        "השתמש ברמות level 1 ו-level 2 בכותרות.\n\n"
        '4. סגנון דיאלוגי: "השאלה שלנו היא...", "נשאל...", "במילים אחרות...", "בואו נבחן...". '
        "שאלות מנחות שמובילות את הקורא דרך הלוגיקה.\n\n"
        "5. דוגמאות בכל מושג מופשט: כל מושג תיאורטי חייב לקבל דוגמה ספציפית. "
        "סמן ב-\"לדוג':\"/\"מקרה:\". אם המרצה נתן — השתמש בה. אם לא — הוסף מהידע שלך.\n\n"
        "6. ציטוט סעיפי חוק וניתוחם: כשמוזכר סעיף — צטט אותו (גם אם המרצה לא קרא — הבא מהידע שלך). "
        "אחרי הציטוט — נתח את רכיביו: יסוד עובדתי (התנהגות / נסיבות / תוצאה), יסוד נפשי. "
        "אם רלוונטי — פרק את הסעיף מילה-מילה.\n\n"
        "7. פסקי דין — ניתוח אקדמי מלא: שם → עובדות → טענות הצדדים → הכרעת בית המשפט → "
        "ההלכה / העיקרון → רלוונטיות. אם המרצה הזכיר רק שם — הרחב מהידע שלך.\n\n"
        "8. שימוש במונחים משפטיים מקצועיים: גם אם המרצה דיבר בשפה פשוטה — תרגם לטרמינולוגיה משפטית "
        "(קוגנטי, פררוגטיבה, השתק שיפוטי, קש\"ס עובדתי, וכו').\n\n"
        "9. הערות הרחבה: כשהמרצה לא פירט משהו, הוסף הערה מסומנת \"💭 הערת הרחבה:\" עם הסבר מהידע "
        "המשפטי שלך. הקורא צריך לדעת מה אמר המרצה ומה אתה הוספת.\n\n"
        "10. הצלבות פנימיות: \"(ראה גם: [נושא])\" כשנושאים מתחברים.\n\n"
        "11. סולם החמרה ויזואלי: כשיש מושגים מדורגים (מחשבה פלילית < רשלנות < אחריות קפידה), "
        "הצג כסקלה מסודרת עם תיאור כל רמה.\n\n"
        "12. טבלאות הבחנה: כשיש מושגים דומים שקל להתבלבל ביניהם (כוונה vs כוונה מיוחדת, מעשה vs מחדל), "
        "צור טבלת השוואה.\n\n"
        "13. סימוני אזהרה: \"⚠️ חשוב לבחינה\" — לנושאים שהמרצה הדגיש (סימוני ⭐ או אינטונציה: "
        "\"שימו לב\", \"זה ייבחן\"). \"❗ פעמים רבות מבלבל בין X ל-Y\". \"💡 טיפ:\".\n\n"
        "14. סיכומי ביניים: בסוף כל סקציית level 1 — הוסף \"📌 סיכום ביניים:\" עם 3-4 נקודות מרכזיות.\n\n"
        "15. הבחנה ויזואלית לפי סוג תוכן בתוך ה-content:\n"
        "- 📜 לשון החוק: ציטוט מדויק\n"
        "- ⚖️ פסיקה: ניתוח אקדמי\n"
        "- 🎓 המרצה הסביר: דוגמה מהשיעור\n"
        "- 💭 הערת הרחבה: ידע משפטי כללי\n\n"
        "16. רגעי חוסר ודאות בתמלול: \"❓ באזור זה התמלול לא היה ברור — בדוק בהקלטה\". "
        "אל תנחש ואל תמציא.\n\n"
        "זיהוי פסקי דין: כל אזכור של X נגד Y, ע\"א, בג\"ץ, רע\"א, \"פרשת X\", \"עניין X\", \"הלכת X\" — "
        "כולם פסקי דין. חלץ עובדות, טיעונים, הכרעה, הלכה.\n\n"
        "זיהוי חוקים: כל \"סעיף X לחוק Y\" — חלץ, צטט, נתח.\n\n"
        "רגעים חשובים: מסומנים ב-⭐ בתמליל. אם אין סימונים — זהה לבד.\n\n"
        "=== פורמט הפלט ===\n\n"
        "החזר JSON בלבד — ללא markdown, ללא טקסט מקדים. מבנה:\n"
        + _JSON_SCHEMA + "\n\n"
        "חשוב: אל תקצר. אל תוותר על איכות. הסטודנט מסתמך על הסיכום הזה במקום השיעור.\n\n"
        + transcript_section
    )


def _chunk_text(text: str) -> list[str]:
    """Split text into overlapping chunks up to MAX_CHUNKS."""
    chunks, start = [], 0
    while start < len(text) and len(chunks) < MAX_CHUNKS:
        chunks.append(text[start : start + CHUNK_SIZE])
        start += CHUNK_SIZE - _CHUNK_OVERLAP
    return chunks


def _extract_chunk_summary(chunk: str, chunk_num: int, total: int, lesson_name: str) -> str:
    """Ask Claude to extract key information from a single transcript chunk."""
    prompt = (
        f'אתה עוזר לימודים. זהה מהקטע הבא (קטע {chunk_num} מתוך {total}) '
        f'משיעור בשם "{lesson_name}":\n\n'
        "- נושאים שנדונו\n"
        "- פסקי דין שהוזכרו (שם תיק + עיקרון)\n"
        "- חוקים וסעיפים שהוזכרו\n"
        "- הגדרות ומושגים\n"
        "- נקודות שהמרצה הדגיש\n\n"
        f"קטע:\n{chunk}"
    )
    return _claude(prompt, max_tokens=2048)


# ─────────────────────────────────────────────────────────────
#  Routes
# ─────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/manifest.json")
def manifest():
    return send_from_directory("static", "manifest.json")


@app.route("/sw.js")
def service_worker():
    return send_from_directory("static", "sw.js", mimetype="application/javascript")


@app.route("/icon-192.svg")
def icon_192():
    return send_from_directory("static", "icon-192.svg", mimetype="image/svg+xml")


@app.route("/icon-512.svg")
def icon_512():
    return send_from_directory("static", "icon-512.svg", mimetype="image/svg+xml")


@app.route("/upload", methods=["POST"])
@limiter.limit("20 per hour")
def upload():
    file = request.files.get("file")
    if not file or not file.filename:
        return jsonify({"error": "לא נבחר קובץ"}), 400

    name = file.filename.lower()
    data = file.read()

    try:
        if name.endswith(".pdf"):
            text = extract_pdf(data)
        elif name.endswith(".pptx"):
            text = extract_pptx(data)
        elif name.endswith(".docx"):
            text = extract_docx(data)
        elif name.endswith(".txt"):
            text = extract_txt(data)
        else:
            return jsonify({"error": "סוג קובץ לא נתמך. קבצים נתמכים: PDF, PPTX, DOCX, TXT"}), 400
    except Exception as exc:
        log.error("Extraction failed: %s", exc)
        return jsonify({"error": f"שגיאה בקריאת הקובץ: {exc}"}), 500

    if not text.strip():
        return jsonify({"error": "לא ניתן לחלץ טקסט מהקובץ"}), 400

    doc_type     = detect_type(text)
    summary_type = request.form.get("summary_type", "full")

    if summary_type == "short":
        prompt = (
            "אתה עוזר לימודים משפטי. סכם את המסמך הבא ב-5–7 נקודות עיקריות בלבד. "
            "קצר, ממוקד, ללא פירוט מיותר. השתמש בפורמט:\n\n"
            "⚡ **סיכום מקוצר**\n\n"
            "1. [נקודה ראשונה]\n"
            "2. [נקודה שנייה]\n"
            "...\n\n"
            "📌 **מסקנה עיקרית:** [משפט אחד]\n\n"
            "המסמך:\n" + text[:12000]
        )
    else:
        prompt = PROMPTS[doc_type] + text[:12000]

    try:
        summary = _claude(prompt)
    except Exception as exc:
        log.error("Claude error: %s", exc)
        return jsonify({"error": f"שגיאה בסיכום: {exc}"}), 500

    token = str(uuid.uuid4())
    if len(_TEXT_CACHE) >= _MAX_CACHE:
        oldest = next(iter(_TEXT_CACHE))
        del _TEXT_CACHE[oldest]
    _TEXT_CACHE[token] = text[:12000]

    icon, label = TYPE_LABELS[doc_type]
    return jsonify({
        "icon": icon, "label": label,
        "doc_type": doc_type,
        "summary": summary, "token": token,
    })


@app.route("/questions", methods=["POST"])
@limiter.limit("20 per hour")
def questions():
    token = (request.json or {}).get("token", "")
    text  = _TEXT_CACHE.get(token)
    if not text:
        return jsonify({"error": "הסשן פג תוקף — אנא העלה את הקובץ מחדש"}), 400

    prompt = (
        "אתה עוזר לימודים משפטי. צור 10 שאלות לבחינה בעברית מהחומר הבא.\n"
        "5 שאלות אמריקאיות (עם 4 תשובות, אחת נכונה) ו-5 שאלות פתוחות.\n"
        "השתמש בפורמט הזה בדיוק:\n\n"
        "📝 שאלות לבחינה\n\n"
        "שאלות אמריקאיות:\n"
        "1. [שאלה]?\n"
        "א. [תשובה]\n"
        "ב. [תשובה]\n"
        "ג. [תשובה]\n"
        "ד. [תשובה]\n"
        "✅ תשובה נכונה: א\n\n"
        "(חזור על הפורמט לשאלות 2-5)\n\n"
        "שאלות פתוחות:\n"
        "1. [שאלה]?\n"
        "💡 תשובה מנחה: [תשובה קצרה]\n\n"
        "(חזור על הפורמט לשאלות 2-5)\n\n"
        "---\n\n" + text
    )
    try:
        return jsonify({"result": _claude(prompt)})
    except Exception as exc:
        log.error("Questions error: %s", exc)
        return jsonify({"error": f"שגיאה ביצירת שאלות: {exc}"}), 500


@app.route("/flashcards", methods=["POST"])
@limiter.limit("20 per hour")
def flashcards():
    token = (request.json or {}).get("token", "")
    text  = _TEXT_CACHE.get(token)
    if not text:
        return jsonify({"error": "הסשן פג תוקף — אנא העלה את הקובץ מחדש"}), 400

    prompt = (
        "אתה עוזר לימודים משפטי. צור 10 כרטיסיות חזרה בעברית מהחומר הבא.\n"
        "כל כרטיסייה: שאלה קצרה מצד אחד, תשובה קצרה מצד שני.\n"
        "השתמש בפורמט הזה בדיוק:\n\n"
        "🃏 כרטיסיות חזרה\n\n"
        "כרטיסייה 1:\n"
        "❓ שאלה: [שאלה קצרה]\n"
        "💡 תשובה: [תשובה קצרה]\n\n"
        "כרטיסייה 2:\n"
        "❓ שאלה: [שאלה קצרה]\n"
        "💡 תשובה: [תשובה קצרה]\n\n"
        "(המשך עד כרטיסייה 10)\n\n"
        "---\n\n" + text
    )
    try:
        return jsonify({"result": _claude(prompt)})
    except Exception as exc:
        log.error("Flashcards error: %s", exc)
        return jsonify({"error": f"שגיאה ביצירת כרטיסיות: {exc}"}), 500


@app.route("/export-docx", methods=["POST"])
def export_docx():
    from docx.shared import RGBColor
    from docx.oxml.ns import qn

    data       = request.json or {}
    summary    = data.get("summary", "")
    filename   = data.get("filename", "סיכום")
    notes      = data.get("notes", "")
    structured = data.get("structured") or {}
    if isinstance(structured, str):
        try:
            structured = json.loads(structured)
        except (json.JSONDecodeError, ValueError):
            structured = {}

    doc = Document()

    def _rtl(paragraph):
        pPr = paragraph._p.get_or_add_pPr()
        bidi = OxmlElement("w:bidi")
        pPr.append(bidi)

    def _heading(text, level):
        h = doc.add_heading(text, level=level)
        _rtl(h)
        return h

    def _cell_write(cell, text, bold=False, white_text=False):
        cell.text = text
        p = cell.paragraphs[0]
        for run in p.runs:
            if bold:
                run.bold = True
            if white_text:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _rtl(p)

    def _cell_bg(cell, hex_color):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"),   "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"),  hex_color)
        tcPr.append(shd)

    def _tbl_rtl(tbl):
        tblPr = tbl._tbl.tblPr
        if tblPr is None:
            tblPr = OxmlElement("w:tblPr")
            tbl._tbl.insert(0, tblPr)
        tblPr.append(OxmlElement("w:bidiVisual"))

    title_para = doc.add_heading(filename, level=0)
    _rtl(title_para)

    if structured.get("sections"):
        # ── Structured path ──────────────────────────────────────
        for sec in structured["sections"]:
            _heading(sec.get("heading", ""), level=min(sec.get("level", 1), 2))
            content = sec.get("content", "")
            if content:
                cp = doc.add_paragraph(content)
                _rtl(cp)

        concepts = structured.get("concepts", [])
        if concepts:
            _heading("💡 טבלת מושגים", level=1)
            tbl = doc.add_table(rows=1, cols=3)
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            hdr = tbl.rows[0].cells
            for cell, title in zip(hdr, ["מושג", "הגדרה", "דוגמה"]):
                _cell_write(cell, title, bold=True, white_text=True)
                _cell_bg(cell, "1E3A5F")
            for c in concepts:
                row = tbl.add_row().cells
                _cell_write(row[0], c.get("term", ""),       bold=True)
                _cell_bg(row[0], "D6E4F7")
                _cell_write(row[1], c.get("definition", ""))
                _cell_bg(row[1], "F5F9FF")
                _cell_write(row[2], c.get("example", ""))
                _cell_bg(row[2], "FEF9EC")
            doc.add_paragraph("")

        case_law = structured.get("case_law", [])
        if case_law:
            _heading("⚖️ טבלת פסיקה", level=1)
            tbl = doc.add_table(rows=1, cols=3)
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            hdr = tbl.rows[0].cells
            for cell, title in zip(hdr, ["שם התיק", "עיקרון", "רלוונטיות"]):
                _cell_write(cell, title, bold=True, white_text=True)
                _cell_bg(cell, "1E3A5F")
            for c in case_law:
                row = tbl.add_row().cells
                _cell_write(row[0], c.get("name", ""),      bold=True)
                _cell_bg(row[0], "EAF0FA")
                _cell_write(row[1], c.get("principle", ""))
                _cell_bg(row[1], "F2F6FC")
                _cell_write(row[2], c.get("relevance", ""))
                _cell_bg(row[2], "F8FAFE")
            doc.add_paragraph("")

        statutes = structured.get("statutes", [])
        if statutes:
            _heading("📜 סעיפי חוק", level=1)
            tbl = doc.add_table(rows=1, cols=3)
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            hdr = tbl.rows[0].cells
            for cell, title in zip(hdr, ["שם החוק", "סעיף", "תוכן"]):
                _cell_write(cell, title, bold=True, white_text=True)
                _cell_bg(cell, "4A235A")
            for s in statutes:
                row = tbl.add_row().cells
                _cell_write(row[0], s.get("law", ""),     bold=True)
                _cell_bg(row[0], "F3E8FA")
                _cell_write(row[1], s.get("section", ""))
                _cell_bg(row[1], "F9F2FD")
                _cell_write(row[2], s.get("content", ""))
                _cell_bg(row[2], "FDF6FF")
            doc.add_paragraph("")

        key_points = structured.get("key_points", [])
        if key_points:
            _heading("📌 נקודות עיקריות", level=1)
            for i, kp in enumerate(key_points[:5], 1):
                kpp = doc.add_paragraph(f"{i}. {kp}")
                _rtl(kpp)

    else:
        # ── Fallback: plain text ─────────────────────────────────
        for line in summary.split("\n"):
            clean = line.strip("*").strip()
            if not clean:
                continue
            p = doc.add_paragraph(clean)
            _rtl(p)

    if notes.strip():
        doc.add_paragraph("")
        h = doc.add_heading("הערות אישיות", level=2)
        _rtl(h)
        for line in notes.split("\n"):
            if line.strip():
                p = doc.add_paragraph(line.strip())
                _rtl(p)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    safe_name = "".join(c for c in filename if c.isalnum() or c in " .-_()") or "summary"
    return send_file(
        buf,
        as_attachment=True,
        download_name=f"{safe_name}.docx",
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@app.route("/summarize-lecture", methods=["POST"])
@limiter.limit("20 per hour")
def summarize_lecture():
    data         = request.json or {}
    lesson_name  = data.get("lesson_name", "שיעור")
    transcript   = data.get("transcript", "").strip()
    duration     = data.get("duration", "")
    summary_type = data.get("summary_type", "full")

    if not transcript:
        return jsonify({"error": "התמלול ריק — ודא שהמיקרופון פעל"}), 400

    today    = date.today().strftime("%d/%m/%Y")
    now_time = datetime.now().strftime("%H:%M")

    # ── Short summary (bullet points, no chunking) ─────────────
    if summary_type == "short":
        prompt = (
            f'אתה עוזר לימודים משפטי. קיבלת תמליל של שיעור בשם: "{lesson_name}"\n'
            f"תאריך: {today}, משך: {duration}\n\n"
            "סכם ב-5–7 נקודות עיקריות בלבד. קצר, ממוקד, ללא פירוט מיותר.\n"
            "השתמש בפורמט:\n\n"
            f"⚡ **סיכום מקוצר — {lesson_name}**\n\n"
            "1. [נקודה ראשונה]\n"
            "2. [נקודה שנייה]\n"
            "...\n\n"
            "📌 **מסקנה עיקרית:** [משפט אחד]\n\n"
            "תמליל השיעור:\n" + transcript[:CHUNK_SIZE]
        )
        try:
            raw = _claude(prompt, max_tokens=1024)
        except Exception as exc:
            log.error("Claude error: %s", exc)
            return jsonify({"error": f"שגיאה בסיכום: {exc}"}), 500
        return jsonify({"summary": raw, "subject": lesson_name, "structured": None})

    # ── Full summary ───────────────────────────────────────────
    try:
        if len(transcript) > CHUNK_SIZE:
            # ── Multi-chunk path: extract then merge ───────────
            chunks = _chunk_text(transcript)
            log.info(
                "Long transcript: %d chars split into %d chunks",
                len(transcript), len(chunks),
            )
            partials = []
            for i, chunk in enumerate(chunks, 1):
                partials.append(_extract_chunk_summary(chunk, i, len(chunks), lesson_name))

            merged_input = "\n\n---\n\n".join(
                f"סיכום קטע {i + 1}:\n{p}" for i, p in enumerate(partials)
            )
            prompt = _full_lecture_prompt(
                lesson_name, today, now_time, duration,
                "סיכומי הביניים:\n" + merged_input,
            )
        else:
            # ── Single-chunk path ──────────────────────────────
            prompt = _full_lecture_prompt(
                lesson_name, today, now_time, duration,
                "תמליל השיעור:\n" + transcript,
            )

        raw        = _claude(prompt, max_tokens=4096)
        js         = raw[raw.find("{") : raw.rfind("}") + 1]
        js         = js.strip().lstrip("```json").lstrip("```").rstrip("```").strip()
        structured = json.loads(js)

    except (json.JSONDecodeError, ValueError):
        log.warning("JSON parse failed — returning plain summary")
        structured = None
    except Exception as exc:
        log.error("Claude error in summarize_lecture: %s", exc)
        return jsonify({"error": f"שגיאה בסיכום: {exc}"}), 500

    if structured is None:
        return jsonify({"summary": raw, "subject": "שיעור", "structured": None})

    subject = structured.get("subject", "שיעור")

    # Build display text
    lines = [
        f"🎓 **נושא:** {subject}",
        f"📅 **תאריך:** {today}   ⏱️ **משך:** {duration}\n",
    ]
    for sec in structured.get("sections", []):
        lines.append(f"\n**{sec.get('heading', '')}**")
        lines.append(sec.get("content", ""))
    if structured.get("concepts"):
        lines.append("\n**💡 מושגים מרכזיים:**")
        for c in structured["concepts"]:
            lines.append(f"• **{c.get('term','')}** — {c.get('definition','')}")
    if structured.get("important_moments"):
        lines.append("\n**⭐ רגעים חשובים:**")
        for m in structured["important_moments"]:
            lines.append(f"• {m}")
    if structured.get("instructor_emphasis"):
        lines.append("\n**📍 נקודות שהמרצה הדגיש:**")
        for e in structured["instructor_emphasis"]:
            lines.append(f"• {e}")
    if structured.get("key_points"):
        lines.append("\n**📌 5 נקודות עיקריות:**")
        for i, kp in enumerate(structured["key_points"][:5], 1):
            lines.append(f"{i}. {kp}")

    return jsonify({
        "summary":    "\n".join(lines),
        "subject":    subject,
        "structured": structured,
    })


@app.route("/export-lecture-docx", methods=["POST"])
def export_lecture_docx():
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn

    data        = request.json or {}
    lesson_name = data.get("lesson_name", "שיעור")
    dt_str      = data.get("date", "")
    duration    = data.get("duration", "")
    subject     = data.get("subject", "")
    structured  = data.get("structured") or {}
    if isinstance(structured, str):
        try:
            structured = json.loads(structured)
        except (json.JSONDecodeError, ValueError):
            structured = {}
    summary_txt = data.get("summary", "")

    doc = Document()

    # ── Helpers ─────────────────────────────────────────────────

    def _rtl(paragraph):
        pPr = paragraph._p.get_or_add_pPr()
        if pPr.find(qn("w:bidi")) is None:
            pPr.append(OxmlElement("w:bidi"))

    def _rtl_right(paragraph):
        _rtl(paragraph)
        paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    def _set_font(run, size_pt=24, bold=False, italic=False, color_hex=None):
        run.bold   = bold
        run.italic = italic
        run.font.size = Pt(size_pt)
        if color_hex:
            run.font.color.rgb = RGBColor(
                int(color_hex[0:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )
        # Set David font for all script types (ascii, hAnsi, complex-script)
        rPr = run._r.get_or_add_rPr()
        rFonts = rPr.find(qn("w:rFonts"))
        if rFonts is None:
            rFonts = OxmlElement("w:rFonts")
            rPr.insert(0, rFonts)
        for attr in ("w:ascii", "w:hAnsi", "w:cs"):
            rFonts.set(qn(attr), "David")

    def _heading_new(text, level):
        p   = doc.add_paragraph()
        run = p.add_run(text)
        _set_font(run, size_pt=(36 if level == 1 else 30), bold=True, color_hex="1F3864")
        _rtl_right(p)
        return p

    def _para(text, size_pt=24):
        p   = doc.add_paragraph()
        run = p.add_run(text)
        _set_font(run, size_pt=size_pt)
        _rtl_right(p)
        return p

    def _cell_bg(cell, hex_color):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd  = OxmlElement("w:shd")
        shd.set(qn("w:val"),   "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"),  hex_color)
        tcPr.append(shd)

    def _cell_write(cell, text, bold=False, white_text=False, italic=False):
        p = cell.paragraphs[0]
        p.clear()
        run = p.add_run(text)
        _set_font(run, size_pt=22, bold=bold, italic=italic)
        if white_text:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _rtl_right(p)

    def _tbl_rtl(tbl):
        tblPr = tbl._tbl.tblPr
        if tblPr is None:
            tblPr = OxmlElement("w:tblPr")
            tbl._tbl.insert(0, tblPr)
        tblPr.append(OxmlElement("w:bidiVisual"))

    # Styled paragraph with colored background and right border
    CONTENT_STYLES = [
        ("📌 לדוג':",        "FEF9EC", "F39C12"),
        ("מקרה:",            "FEF9EC", "F39C12"),
        ("📜 לשון הסעיף:",   "EAF0FA", "2E5090"),
        ("⚖️ פסק הדין:",     "D1FAE5", "059669"),
        ("💭 הערת הרחבה:",   "EDE9FE", "8B5CF6"),
        ("⚠️ חשוב לבחינה:", "FEE2E2", "DC2626"),
        ("📌 סיכום ביניים:", "D1FAE5", "059669"),
    ]

    def _styled_para(text, bg_hex, border_hex):
        p   = doc.add_paragraph()
        run = p.add_run(text)
        _set_font(run, size_pt=24)
        pPr = p._p.get_or_add_pPr()
        _rtl(p)
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"),   "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"),  bg_hex)
        pPr.append(shd)
        pBdr = OxmlElement("w:pBdr")
        rb   = OxmlElement("w:right")
        rb.set(qn("w:val"),   "single")
        rb.set(qn("w:sz"),    "24")
        rb.set(qn("w:space"), "4")
        rb.set(qn("w:color"), border_hex)
        pBdr.append(rb)
        pPr.append(pBdr)
        ind = OxmlElement("w:ind")
        ind.set(qn("w:right"), "180")
        pPr.append(ind)
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        return p

    def _render_content(text):
        for line in text.split("\n"):
            if not line.strip():
                continue
            matched = False
            for prefix, bg, border in CONTENT_STYLES:
                if line.startswith(prefix):
                    _styled_para(line, bg, border)
                    matched = True
                    break
            if not matched:
                _para(line)

    # ── Cover page ──────────────────────────────────────────────
    cover_p = doc.add_paragraph()
    cover_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cover_run = cover_p.add_run(lesson_name)
    _set_font(cover_run, size_pt=60, bold=True, color_hex="1F3864")
    _rtl(cover_p)

    if subject:
        sub_p   = doc.add_paragraph()
        sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = sub_p.add_run(subject)
        _set_font(sub_run, size_pt=36, color_hex="2E5090")
        _rtl(sub_p)

    meta_parts = []
    if dt_str:   meta_parts.append(f"📅 תאריך: {dt_str}")
    if duration: meta_parts.append(f"⏱️ משך: {duration}")
    if meta_parts:
        meta_p = doc.add_paragraph("   ".join(meta_parts))
        meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        _rtl(meta_p)

    doc.add_page_break()

    # ── Body sections with inline styled boxes ───────────────────
    for sec in structured.get("sections", []):
        _heading_new(sec.get("heading", ""), level=min(sec.get("level", 1), 2))
        content = sec.get("content", "")
        if content:
            _render_content(content)

    # ── Comparison tables ────────────────────────────────────────
    comparison_tables = structured.get("comparison_tables", [])
    if comparison_tables:
        _heading_new("📊 טבלאות השוואה", 1)
        for ct in comparison_tables:
            _heading_new(ct.get("title", ""), 2)
            columns   = ct.get("columns", [])
            rows_data = ct.get("rows", [])
            if not columns:
                continue
            tbl = doc.add_table(rows=1, cols=len(columns))
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            for j, col_name in enumerate(columns):
                _cell_write(tbl.rows[0].cells[j], col_name, bold=True, white_text=True)
                _cell_bg(tbl.rows[0].cells[j], "1F3864")
            for i, row_data in enumerate(rows_data):
                tbl_row = tbl.add_row().cells
                bg = "F2F6FC" if i % 2 == 0 else "FFFFFF"
                for j, val in enumerate(row_data):
                    _cell_write(tbl_row[j], str(val) if val else "")
                    _cell_bg(tbl_row[j], bg)
            doc.add_paragraph("")

    # ── Scales ───────────────────────────────────────────────────
    scales = structured.get("scales", [])
    if scales:
        _heading_new("📊 סולמות", 1)
        for scale in scales:
            _heading_new(scale.get("title", ""), 2)
            levels_list = scale.get("levels", [])
            if not levels_list:
                continue
            n   = len(levels_list)
            tbl = doc.add_table(rows=n, cols=2)
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            for i, lv in enumerate(levels_list):
                frac  = i / max(n - 1, 1)
                r_val = int(0x1F + (0x44 - 0x1F) * frac)
                g_val = int(0x38 + (0x72 - 0x38) * frac)
                b_val = int(0x64 + (0xC4 - 0x64) * frac)
                row   = tbl.rows[i].cells
                _cell_write(row[0], lv.get("name", ""), bold=True, white_text=True)
                _cell_bg(row[0], f"{r_val:02X}{g_val:02X}{b_val:02X}")
                _cell_write(row[1], lv.get("description", ""))
                _cell_bg(row[1], "EAF0FA")
            doc.add_paragraph("")

    # ── Important moments ────────────────────────────────────────
    moments = structured.get("important_moments", [])
    if moments:
        _heading_new("⭐ רגעים חשובים", 1)
        for m in moments:
            p   = doc.add_paragraph()
            run = p.add_run(m)
            _set_font(run, size_pt=24, bold=True, color_hex="B88600")
            _rtl_right(p)

    # ── Concepts table ───────────────────────────────────────────
    concepts = structured.get("concepts", [])
    if concepts:
        _heading_new("💡 טבלת מושגים", 1)
        tbl = doc.add_table(rows=1, cols=3)
        tbl.style = "Table Grid"
        _tbl_rtl(tbl)
        for cell, title in zip(tbl.rows[0].cells, ["מושג", "הגדרה", "דוגמה"]):
            _cell_write(cell, title, bold=True, white_text=True)
            _cell_bg(cell, "1E3A5F")
        for c in concepts:
            row = tbl.add_row().cells
            _cell_write(row[0], c.get("term", ""),       bold=True)
            _cell_bg(row[0], "D6E4F7")
            _cell_write(row[1], c.get("definition", ""))
            _cell_bg(row[1], "F5F9FF")
            _cell_write(row[2], c.get("example", ""))
            _cell_bg(row[2], "FEF9EC")
        doc.add_paragraph("")

    # ── Case law — vertical per case ─────────────────────────────
    case_law = structured.get("case_law", [])
    if case_law:
        _heading_new("⚖️ פסיקה", 1)
        fields = [
            ("עובדות",              "facts"),
            ("טענות הצדדים",        "arguments"),
            ("הכרעת בית המשפט",    "ruling"),
            ("ההלכה / העיקרון",    "principle"),
            ("רלוונטיות",           "relevance"),
        ]
        for c in case_law:
            tbl = doc.add_table(rows=6, cols=2)
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            title_cells = tbl.rows[0].cells
            title_cells[0].merge(title_cells[1])
            _cell_write(title_cells[0], c.get("name", ""), bold=True, white_text=True)
            _cell_bg(title_cells[0], "1F3864")
            for i, (label, key) in enumerate(fields):
                row = tbl.rows[i + 1].cells
                _cell_write(row[0], label, bold=True)
                _cell_bg(row[0], "EAF0FA")
                _cell_write(row[1], c.get(key, ""))
                _cell_bg(row[1], "F8FAFE")
            doc.add_paragraph("")

    # ── Statutes ─────────────────────────────────────────────────
    statutes = structured.get("statutes", [])
    if statutes:
        _heading_new("📜 סעיפי חוק", 1)
        for s in statutes:
            tbl = doc.add_table(rows=3, cols=1)
            tbl.style = "Table Grid"
            _tbl_rtl(tbl)
            _cell_write(tbl.rows[0].cells[0],
                        f"{s.get('law','')} — {s.get('section','')}", bold=True, white_text=True)
            _cell_bg(tbl.rows[0].cells[0], "1F3864")
            _cell_write(tbl.rows[1].cells[0],
                        "📜 לשון הסעיף: " + s.get("text", ""), italic=True)
            _cell_bg(tbl.rows[1].cells[0], "EAF0FA")
            _cell_write(tbl.rows[2].cells[0],
                        "🔍 ניתוח: " + s.get("analysis", ""))
            _cell_bg(tbl.rows[2].cells[0], "F2F6FC")
            doc.add_paragraph("")

    # ── Exam warnings ────────────────────────────────────────────
    exam_warnings = structured.get("exam_warnings", [])
    if exam_warnings:
        _heading_new("⚠️ אזהרות לבחינה", 1)
        for w in exam_warnings:
            _styled_para("⚠️ " + w, "FEE2E2", "DC2626")

    # ── Related topics ───────────────────────────────────────────
    related = structured.get("related_topics", "")
    if related:
        _heading_new("🔗 קשרים לנושאים אחרים", 1)
        _para(related)

    # ── Instructor emphasis ──────────────────────────────────────
    emphasis = structured.get("instructor_emphasis", [])
    if emphasis:
        _heading_new("📍 נקודות לבדיקה", 1)
        for e in emphasis:
            _para("• " + e)

    # ── Key points summary ───────────────────────────────────────
    key_points = structured.get("key_points", [])
    if key_points:
        _heading_new("📌 5 נקודות עיקריות מהשיעור", 1)
        for i, kp in enumerate(key_points[:5], 1):
            _para(f"{i}. {kp}")

    # ── Transcript unclear zones ─────────────────────────────────
    unclear = structured.get("transcript_unclear_zones", [])
    if unclear:
        _heading_new("❓ אזורי תמלול לא ברורים", 1)
        for zone in unclear:
            _styled_para("❓ " + zone, "F3F4F6", "9CA3AF")

    # ── Fallback: plain summary ──────────────────────────────────
    if not structured.get("sections") and summary_txt:
        _heading_new("סיכום", 1)
        for line in summary_txt.split("\n"):
            clean = line.strip("*").strip()
            if clean:
                _para(clean)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    safe = "".join(c for c in lesson_name if c.isalnum() or c in " .-_()") or "שיעור"
    return send_file(
        buf, as_attachment=True,
        download_name=f"{safe}.docx",
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@app.route("/concept-map", methods=["POST"])
@limiter.limit("20 per hour")
def concept_map():
    token = (request.json or {}).get("token", "")
    text  = _TEXT_CACHE.get(token)
    if not text:
        return jsonify({"error": "הסשן פג תוקף — אנא העלה את הקובץ מחדש"}), 400

    prompt = (
        "אתה עוזר לימודים משפטי. נתח את המסמך הבא וזהה:\n"
        "- פסקי דין שהוזכרו (type: verdict)\n"
        "- חוקים וסעיפים (type: law)\n"
        "- עקרונות משפטיים (type: principle)\n"
        "- מושגים משפטיים חשובים (type: concept)\n"
        "- קשרים בין האלמנטים הללו\n\n"
        "החזר JSON בלבד — ללא markdown, ללא טקסט לפני או אחרי הסוגריים:\n"
        '{\n'
        '  "nodes": [\n'
        '    {"id":"1","label":"שם קצר (עד 35 תווים)","type":"verdict|law|principle|concept"}\n'
        '  ],\n'
        '  "edges": [\n'
        '    {"from":"1","to":"2","label":"קשר קצר (עד 20 תווים)"}\n'
        '  ]\n'
        '}\n\n'
        "כלול 6–15 צמתים ו-5–15 קשרים. אם אין פסקי דין — השמט. label לכל צומת חובה.\n\n"
        "המסמך:\n" + text[:9000]
    )

    try:
        raw  = _claude(prompt, max_tokens=2048)
        js   = raw[raw.find("{") : raw.rfind("}") + 1]
        data = json.loads(js)
        # Basic validation
        if not data.get("nodes"):
            return jsonify({"error": "לא זוהו אלמנטים משפטיים במסמך"}), 400
    except (json.JSONDecodeError, ValueError) as exc:
        log.error("Concept-map JSON error: %s", exc)
        return jsonify({"error": "שגיאה בניתוח המסמך — נסה שוב"}), 500

    return jsonify(data)


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 10000))
    app.run(host='0.0.0.0', port=port, debug=False)
